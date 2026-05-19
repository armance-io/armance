from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Union

from markdown_it import MarkdownIt

# Third-party renderers
try:
    from docx import Document
    HAS_DOCX = True
except Exception:  # pragma: no cover
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except Exception:  # pragma: no cover
    HAS_PPTX = False

try:
    from weasyprint import HTML
    HAS_WEASYPRINT = True
except Exception:  # pragma: no cover
    HAS_WEASYPRINT = False


class DeliverableError(RuntimeError):
    """Raised when a deliverable cannot be rendered (e.g. missing dependency)."""


@dataclass
class Section:
    heading: str
    heading_level: int
    body: str
    lists: List[str]


@dataclass
class ReportTree:
    title: str
    sections: List[Section]


def parse_report(md_text: str) -> ReportTree:
    """Parse a markdown report into a structured ReportTree using raw tokens.

    Rules:
    - H1 → title only, no section created
    - H2 → creates a section
    - H3+ heading → list item in parent H2; marks start of H3 sub-section
    - H2-level bullets → section.lists only
    - H3-sub bullets/paragraphs → section.body only
    - H2 paragraphs → section.body
    """
    md = MarkdownIt("commonmark")
    tokens = md.parse(md_text)

    title = ""
    sections: List[Section] = []
    current_heading: Optional[str] = None
    current_heading_level: int = 0
    current_body: List[str] = []
    current_lists: List[str] = []

    in_heading = False
    current_tag_level: int = 0
    in_paragraph = False
    in_list_item = False
    in_h3_subsection = False

    for tok in tokens:
        if tok.type == "heading_open":
            in_heading = True
            current_tag_level = int(tok.tag[1:])

        elif tok.type == "heading_close":
            in_heading = False

        elif tok.type == "paragraph_open":
            in_paragraph = True

        elif tok.type == "paragraph_close":
            in_paragraph = False

        elif tok.type == "list_item_open":
            in_list_item = True

        elif tok.type == "list_item_close":
            in_list_item = False

        elif tok.type == "inline":
            content = tok.content
            if not content.strip():
                continue

            if in_heading:
                if current_tag_level == 1:
                    if not title:
                        title = content
                    # H1 never becomes a section
                elif current_tag_level == 2:
                    if current_heading is not None:
                        sections.append(Section(
                            heading=current_heading,
                            heading_level=current_heading_level,
                            body="\n".join(current_body),
                            lists=current_lists[:],
                        ))
                    current_heading = content
                    current_heading_level = 2
                    current_body = []
                    current_lists = []
                    in_h3_subsection = False
                else:
                    # H3+ → list item in parent H2, enter H3 sub-section
                    if current_heading is not None:
                        current_lists.append(content)
                    in_h3_subsection = True

            elif in_list_item and in_paragraph and current_heading is not None:
                if in_h3_subsection:
                    current_body.append(content)
                else:
                    current_lists.append(content)

            elif in_paragraph and not in_list_item and current_heading is not None:
                current_body.append(content)

    if current_heading is not None:
        sections.append(Section(
            heading=current_heading,
            heading_level=current_heading_level,
            body="\n".join(current_body),
            lists=current_lists[:],
        ))

    return ReportTree(title=title or "Report", sections=sections)


def render_docx(
    tree: ReportTree,
    out_path: Union[Path, str],
    template: Optional[Union[Path, str]] = None,
) -> None:
    """
    Render a ReportTree to a DOCX file.
    
    Creates a document with:
    - Title page (title only)
    - Table of contents (simple text list of sections)
    - Sections as Heading 1/2 styles
    - Body paragraphs
    - Bullet lists preserved
    """
    if not HAS_DOCX:
        raise RuntimeError(
            "python-docx not available. Install with: pip install python-docx"
        )

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    doc = Document()

    # Default styles
    doc.add_heading(tree.title, level=1)

    # Table of contents (simple text list)
    if tree.sections:
        doc.add_heading("Table of Contents", level=2)
        for section in tree.sections:
            doc.add_paragraph(section.heading)

    # Sections
    for section in tree.sections:
        doc.add_heading(section.heading, level=section.heading_level)
        
        # Body paragraphs
        if section.body:
            for para in section.body.split("\n"):
                if para.strip():
                    doc.add_paragraph(para)
        
        # Bullet lists
        for lst_item in section.lists:
            p = doc.add_paragraph(style="List Bullet")
            p.add_run(lst_item)

    doc.save(out_path)


def render_pptx(
    tree: ReportTree,
    out_path: Union[Path, str],
    template: Optional[Union[Path, str]] = None,
) -> None:
    """
    Render a ReportTree to a PPTX file.
    
    Creates a presentation with:
    - Title slide (title only)
    - Summary slide (first paragraph of first section)
    - One slide per H1 section
    - Bullet points (max 6 per slide, auto-split if more)
    """
    if not HAS_PPTX:
        raise RuntimeError(
            "python-pptx not available. Install with: pip install python-pptx"
        )

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = tree.title
    subtitle.text = "Generated by Armance"

    # Summary slide (first paragraph of first section)
    if tree.sections and tree.sections[0].body:
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Summary"
        
        # Add first paragraph
        p = content.text_frame.add_paragraph()
        p.text = tree.sections[0].body.split("\n")[0]
        p.level = 0

    # Content slides (one per section)
    for section in tree.sections:
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        slide_title = slide.shapes.title
        content = slide.placeholders[1]

        slide_title.text = section.heading

        tf = content.text_frame
        tf.clear()

        # Body paragraphs
        for para in section.body.split("\n"):
            if para.strip():
                p = tf.add_paragraph()
                p.text = para
                p.level = 0

        # Bullet lists - max 6 per slide, auto-split into new slides
        items = list(section.lists)
        for chunk_start in range(0, max(1, len(items)), 6):
            chunk = items[chunk_start:chunk_start + 6]
            if chunk_start > 0:
                # Overflow slide
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = section.heading + " (cont.)"
                content = slide.placeholders[1]
                tf = content.text_frame
                tf.clear()
            for lst_item in chunk:
                p = tf.add_paragraph()
                p.text = lst_item
                p.level = 1  # Bullet level

    prs.save(out_path)


def render_pdf(
    tree: ReportTree,
    out_path: Union[Path, str],
) -> None:
    """
    Render a ReportTree to a PDF file using WeasyPrint.
    
    Converts markdown -> HTML (via markdown-it) -> CSS-styled PDF.
    Uses default template from src/armance/templates/pdf_default.css
    """
    if not HAS_WEASYPRINT:
        raise DeliverableError(
            "WeasyPrint not available. "
            "Install system deps (libgobject-2.0-0, libcairo2, libpango-1.0-0) then: "
            "pip install weasyprint"
        )

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    # Build HTML from markdown
    md = MarkdownIt("commonmark")
    html = md.render(
        f"# {tree.title}\n\n" + "\n\n".join(
            [f"{'#' * section.heading_level} {section.heading}\n\n{section.body}" for section in tree.sections]
        )
    )

    # Load default CSS
    css_path = Path(__file__).parent / "templates" / "pdf_default.css"
    if css_path.exists():
        css = css_path.read_text(encoding="utf-8")
    else:
        css = """
        body { font-family: Arial, sans-serif; font-size: 12pt; line-height: 1.5; }
        h1 { font-size: 24pt; font-weight: bold; margin-top: 24pt; margin-bottom: 12pt; }
        h2 { font-size: 18pt; font-weight: bold; margin-top: 20pt; margin-bottom: 8pt; }
        p { margin-top: 0; margin-bottom: 8pt; }
        ul { margin-top: 0; margin-bottom: 8pt; padding-left: 24pt; }
        li { margin-top: 4pt; margin-bottom: 4pt; }
        """

    # Render to PDF
    from weasyprint import CSS
    HTML(string=html, base_url=str(out_path.parent)).write_pdf(
        str(out_path),
        stylesheets=[CSS(string=css)],
        presentational_hints=True,
    )
