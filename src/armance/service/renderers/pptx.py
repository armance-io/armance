"""PPTX renderer — one slide per top-level heading.

Requires python-pptx (optional [render] extra). Falls back gracefully.

Spec: docs/spec/22_circular_outputs.md § Supported formats (pptx)
"""
from __future__ import annotations

import logging
from pathlib import Path

from .base import Deliverable, RenderResult

logger = logging.getLogger(__name__)


def _parse_slides(content: str) -> list[tuple[str, list[str]]]:
    """Return list of (title, [bullet lines]) from markdown headings."""
    slides: list[tuple[str, list[str]]] = []
    current_title = ""
    current_bullets: list[str] = []
    for line in content.splitlines():
        if line.startswith("# "):
            if current_title:
                slides.append((current_title, current_bullets))
            current_title = line[2:].strip()
            current_bullets = []
        elif line.startswith("## "):
            current_bullets.append(line[3:].strip())
        elif line.strip() and not line.startswith("#"):
            current_bullets.append(line.strip())
    if current_title:
        slides.append((current_title, current_bullets))
    if not slides:
        slides = [("Output", content.splitlines()[:10])]
    return slides


class PptxRenderer:
    """Renders deliverables into a PPTX presentation."""

    format = "pptx"

    async def render(
        self,
        deliverables: list[Deliverable],
        template: Path | None,
        options: dict,
        output_path: Path,
    ) -> RenderResult:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return RenderResult(
                output_path=output_path,
                warnings=["python-pptx not installed; install with: pip install armance[render]"],
                error="python-pptx not installed",
            )

        combined = "\n\n".join(d.content for d in deliverables)
        slides_data = _parse_slides(combined)

        prs = Presentation(str(template)) if template and template.exists() else Presentation()
        # Clear template slides if using one
        if template and template.exists() and prs.slides:
            pass  # keep template structure

        slide_layout = prs.slide_layouts[1]  # title_and_content
        for title, bullets in slides_data:
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for bullet in bullets[:10]:
                p = tf.add_paragraph()
                p.text = bullet

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        size = output_path.stat().st_size
        return RenderResult(
            output_path=output_path,
            bytes_written=size,
            pages_or_slides=len(slides_data),
        )
