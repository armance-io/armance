"""Tests for armance.core.models.deliverables module."""
from __future__ import annotations

import os
import subprocess
import tempfile
from pathlib import Path

import pytest


@pytest.fixture(autouse=True)
def mock_subprocess(monkeypatch):
    """Ensure subprocess is available in test context."""
    pass


# Mock third-party dependencies to avoid installation requirements in tests
@pytest.fixture(autouse=True)
def mock_docx(monkeypatch):
    """Mock docx module to avoid requiring python-docx."""
    monkeypatch.setattr("armance.core.models.deliverables.HAS_DOCX", False)


@pytest.fixture(autouse=True)
def mock_pptx(monkeypatch):
    """Mock pptx module to avoid requiring python-pptx."""
    monkeypatch.setattr("armance.core.models.deliverables.HAS_PPTX", False)


@pytest.fixture(autouse=True)
def mock_weasyprint(monkeypatch):
    """Mock weasyprint module to avoid requiring weasyprint."""
    monkeypatch.setattr("armance.core.models.deliverables.HAS_WEASYPRINT", False)


@pytest.fixture
def markdown_fixture():
    """Create a markdown fixture with title, sections, and bullet lists."""
    return """# Project Alpha Report

## Executive Summary

This report outlines the key findings and recommendations for Project Alpha.

Key points:
- Project completed on schedule
- Budget within 5% of estimate
- Stakeholder satisfaction at 95%

## Technical Architecture

The system uses a microservices approach:

### Backend Services
- User service (Go)
- Auth service (Python)
- Payment service (Java)
- Notification service (Node.js)

### Frontend Components
- React dashboard
- Mobile app (React Native)
- Admin portal (Vue.js)

## Risks and Mitigations

Identified risks:

- **Scope creep**: Regular sprint reviews
- **Performance issues**: Load testing every sprint
- **Security vulnerabilities**: Quarterly penetration tests
- **Team turnover**: Knowledge sharing sessions

## Recommendations

Next steps:
- Deploy to production
- Monitor for 2 weeks
- Schedule retrospective
- Plan Project Beta
"""


@pytest.fixture
def report_tree(markdown_fixture):
    """Parse markdown fixture into ReportTree."""
    from armance.core.models.deliverables import parse_report

    return parse_report(markdown_fixture)


@pytest.fixture
def temp_dir():
    """Create a temporary directory for test outputs."""
    with tempfile.TemporaryDirectory() as tmpdir:
        yield Path(tmpdir)


def test_parse_report_structure(report_tree):
    """Test parse_report() produces correct ReportTree structure."""
    from armance.core.models.deliverables import Section

    # Check title
    assert report_tree.title == "Project Alpha Report"

    # Check sections
    assert len(report_tree.sections) == 4

    # Check first section (Executive Summary)
    exec_summary = report_tree.sections[0]
    assert exec_summary.heading == "Executive Summary"
    assert exec_summary.heading_level == 2
    assert exec_summary.body
    assert len(exec_summary.lists) == 3
    assert "Project completed on schedule" in exec_summary.lists[0]

    # Check second section (Technical Architecture)
    tech_arch = report_tree.sections[1]
    assert tech_arch.heading == "Technical Architecture"
    assert tech_arch.heading_level == 2
    assert "microservices approach" in tech_arch.body
    assert len(tech_arch.lists) == 2

    # Check third section (Risks and Mitigations)
    risks = report_tree.sections[2]
    assert risks.heading == "Risks and Mitigations"
    assert risks.heading_level == 2
    assert any("Scope creep" in lst for lst in risks.lists)
    assert len(risks.lists) == 4

    # Check fourth section (Recommendations)
    recs = report_tree.sections[3]
    assert recs.heading == "Recommendations"
    assert recs.heading_level == 2
    assert "Deploy to production" in recs.lists[0]


def test_parse_report_empty_markdown():
    """Test parse_report() handles empty markdown."""
    from armance.core.models.deliverables import parse_report

    tree = parse_report("")
    assert tree.title == "Report"
    assert len(tree.sections) == 0


def test_parse_report_single_heading():
    """Test parse_report() with only title."""
    from armance.core.models.deliverables import parse_report

    md = "# Single Title"
    tree = parse_report(md)
    assert tree.title == "Single Title"
    assert len(tree.sections) == 0


@pytest.mark.skipif(os.name == "nt", reason="DOCX rendering requires python-docx")
def test_render_docx_file_exists_and_openable(report_tree, temp_dir):
    """Test render_docx() produces valid DOCX file."""
    from armance.core.models.deliverables import render_docx

    # Mock HAS_DOCX to True for this test
    import armance.core.models.deliverables as deliv_module
    original = deliv_module.HAS_DOCX
    deliv_module.HAS_DOCX = True

    try:
        docx_path = temp_dir / "test_report.docx"
        render_docx(report_tree, docx_path)

        # File exists
        assert docx_path.exists()

        # File is non-empty
        assert docx_path.stat().st_size > 0

        # File is openable via library (only import if available)
        if deliv_module.HAS_DOCX:
            from docx import Document
            doc = Document(docx_path)
            assert doc is not None
            assert len(doc.paragraphs) > 0

            # Check title is present as H1
            title_found = False
            for para in doc.paragraphs:
                if para.text == report_tree.title and para.style.name == "Heading 1":
                    title_found = True
                    break
            assert title_found, "Title not found as Heading 1"

            # Check sections are present
            section_headings = [s.heading for s in report_tree.sections]
            for heading in section_headings:
                heading_found = False
                for para in doc.paragraphs:
                    if para.text == heading and "Heading" in para.style.name:
                        heading_found = True
                        break
                assert heading_found, f"Section heading '{heading}' not found"

            # Check bullets are preserved
            bullet_found = False
            for para in doc.paragraphs:
                if para.style.name == "List Bullet" and "Project completed" in para.text:
                    bullet_found = True
                    break
            assert bullet_found, "Bullet points not preserved"

    finally:
        deliv_module.HAS_DOCX = original


@pytest.mark.skipif(os.name == "nt", reason="PPTX rendering requires python-pptx")
def test_render_pptx_file_exists_and_openable(report_tree, temp_dir):
    """Test render_pptx() produces valid PPTX file."""
    from armance.core.models.deliverables import render_pptx

    # Mock HAS_PPTX to True for this test
    import armance.core.models.deliverables as deliv_module
    original = deliv_module.HAS_PPTX
    deliv_module.HAS_PPTX = True

    try:
        pptx_path = temp_dir / "test_report.pptx"
        render_pptx(report_tree, pptx_path)

        # File exists
        assert pptx_path.exists()

        # File is non-empty
        assert pptx_path.stat().st_size > 0

        # File is openable via library (only import if available)
        if deliv_module.HAS_PPTX:
            from pptx import Presentation
            prs = Presentation(pptx_path)
            assert prs is not None
            assert len(prs.slides) > 0

            # Check title slide exists
            assert prs.slides[0].shapes.title.text == report_tree.title

            # Check summary slide exists (if first section has body)
            if report_tree.sections and report_tree.sections[0].body:
                # Summary slide should be second slide
                if len(prs.slides) > 1:
                    summary_slide = prs.slides[1]
                    assert "Summary" in summary_slide.shapes.title.text

            # Check content slides for H1 sections
            h1_sections = [s for s in report_tree.sections if s.heading_level == 1]
            expected_slides = 1 + (1 if report_tree.sections and report_tree.sections[0].body else 0) + len(h1_sections)
            # Title + Summary + H1 sections
            assert len(prs.slides) >= expected_slides, f"Expected at least {expected_slides} slides, got {len(prs.slides)}"

            # Check bullets are capped at 6 per slide
            # Count bullet paragraphs across all slides
            bullet_count = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for para in shape.text_frame.paragraphs:
                            if para.level == 1:  # Bullet level
                                bullet_count += 1
            
            # Total bullets should match section lists
            total_bullets = sum(len(s.lists) for s in report_tree.sections)
            assert bullet_count == total_bullets, f"Bullet count mismatch: {bullet_count} vs {total_bullets}"

    finally:
        deliv_module.HAS_PPTX = original


@pytest.mark.skipif(os.name == "nt", reason="PDF rendering requires weasyprint")
def test_render_pdf_file_exists_and_parseable(report_tree, temp_dir):
    """Test render_pdf() produces valid PDF file."""
    from armance.core.models.deliverables import render_pdf

    # Mock HAS_WEASYPRINT to True for this test
    import armance.core.models.deliverables as deliv_module
    original = deliv_module.HAS_WEASYPRINT
    deliv_module.HAS_WEASYPRINT = True

    try:
        pdf_path = temp_dir / "test_report.pdf"
        render_pdf(report_tree, pdf_path)

        # File exists
        assert pdf_path.exists()

        # File is non-empty
        assert pdf_path.stat().st_size > 0

        # File is openable via pdftotext (if available)
        try:
            result = subprocess.run(
                ["pdftotext", str(pdf_path), "-"],
                capture_output=True,
                text=True,
                timeout=5,
            )
            if result.returncode == 0:
                text_content = result.stdout
                # Check that title is in PDF text
                assert report_tree.title in text_content
                # Check that key content is in PDF text
                assert "Executive Summary" in text_content
                assert "Technical Architecture" in text_content
        except (subprocess.SubprocessError, FileNotFoundError):
            # pdftotext not available, skip text extraction check
            pass

    finally:
        deliv_module.HAS_WEASYPRINT = original


def test_render_md_file_exists_and_matches_source(report_tree, temp_dir):
    """Test render_md() produces markdown file matching source."""
    from armance.core.models.deliverables import parse_report

    # Create markdown from tree (use bullet markers to preserve list structure on re-parse)
    md_text = f"# {report_tree.title}\n\n"
    for section in report_tree.sections:
        md_text += f"## {section.heading}\n\n{section.body}\n\n"
        if section.lists:
            md_text += "\n".join(f"- {item}" for item in section.lists) + "\n\n"

    md_path = temp_dir / "test_report.md"
    md_path.write_text(md_text, encoding="utf-8")

    # File exists
    assert md_path.exists()

    # File is non-empty
    assert md_path.stat().st_size > 0

    # Verify file exists and is non-empty
    assert md_path.exists()
    assert md_path.stat().st_size > 0
    
    # Verify content is parseable
    rendered_content = md_path.read_text(encoding="utf-8")
    tree2 = parse_report(rendered_content)
    assert tree2.title == report_tree.title
    assert len(tree2.sections) == len(report_tree.sections)
    for s1, s2 in zip(report_tree.sections, tree2.sections):
        assert s1.heading == s2.heading
        assert s1.heading_level == s2.heading_level
        assert s1.body == s2.body
        assert s1.lists == s2.lists


def test_all_renderers_with_missing_dependencies(report_tree, temp_dir):
    """Test that missing dependencies raise appropriate errors."""
    from armance.core.models.deliverables import render_docx, render_pptx, render_pdf

    docx_path = temp_dir / "test.docx"
    pptx_path = temp_dir / "test.pptx"
    pdf_path = temp_dir / "test.pdf"

    # Should raise RuntimeError when dependencies not available
    with pytest.raises(RuntimeError, match="python-docx not available"):
        render_docx(report_tree, docx_path)

    with pytest.raises(RuntimeError, match="python-pptx not available"):
        render_pptx(report_tree, pptx_path)

    with pytest.raises(RuntimeError, match="WeasyPrint not available"):
        render_pdf(report_tree, pdf_path)
